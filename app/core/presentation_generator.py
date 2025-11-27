"""
GENERADOR DE PRESENTACIONES POWERPOINT
Crea slides profesionales basados en contenido del sílabo
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import json
import os
import requests
from io import BytesIO
from typing import Dict, List
from pathlib import Path
from openai import OpenAI

class PresentationGenerator:
    """Genera presentaciones de PowerPoint profesionales"""
    
    def __init__(self, use_images=True):
        self.output_dir = Path("presentations")
        self.output_dir.mkdir(exist_ok=True)
        self.use_images = use_images
        self.openai_client = None
        
        # Inicializar cliente de OpenAI si hay API key
        if use_images and os.getenv("OPENAI_API_KEY"):
            try:
                self.openai_client = OpenAI(api_key=os.getenv("OPENAI_API_KEY"))
                print("   🎨 DALL-E activado para generar imágenes")
            except Exception as e:
                print(f"   ⚠️  DALL-E no disponible: {str(e)}")
                self.use_images = False
    
    def create_presentation(self, title: str, subtitle: str, slides_data: List[Dict]) -> str:
        """
        Crea una presentación completa
        
        Args:
            title: Título principal
            subtitle: Subtítulo (ej: "Unidad 2 - Cálculo")
            slides_data: Lista de diccionarios con datos de cada slide
            
        Returns:
            Path al archivo .pptx generado
        """
        # Crear presentación
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        
        # Slide 1: Título
        self._add_title_slide(prs, title, subtitle)
        
        # Slides de contenido
        for slide_data in slides_data:
            slide_type = slide_data.get('type', 'content')
            
            if slide_type == 'content':
                self._add_content_slide(prs, slide_data)
            elif slide_type == 'two_column':
                self._add_two_column_slide(prs, slide_data)
            elif slide_type == 'summary':
                self._add_summary_slide(prs, slide_data)
        
        # Guardar archivo
        print(f"   🏷️  Título original: {title}")
        filename = self._sanitize_filename(title)
        print(f"   🏷️  Filename sanitizado: {filename}")
        filepath = self.output_dir / f"{filename}.pptx"
        print(f"   💾 Guardando en: {filepath}")
        prs.save(str(filepath))
        
        return str(filepath)
    
    def _generate_image(self, prompt: str) -> BytesIO:
        """
        Genera una imagen con DALL-E 2 basada en el prompt
        
        Args:
            prompt: Descripción de la imagen a generar
            
        Returns:
            BytesIO con los datos de la imagen
        """
        if not self.openai_client:
            return None
        
        try:
            # Traducir conceptos clave al inglés para mejor resultado
            prompt_translations = {
                # Historia
                "Era del Guano": "Guano Era Peru",
                "Guerra del Pacífico": "War of the Pacific",
                "Guerra con Chile": "War with Chile",
                "Ramón Castilla": "Ramon Castilla Peru president",
                "Prosperidad Falaz": "False Prosperity Peru",
                
                # Machine Learning
                "Machine Learning Supervisado": "Supervised Machine Learning",
                "Aprendizaje Supervisado": "Supervised Learning",
                "Redes Neuronales": "Neural Networks",
                "Algoritmos": "Algorithms",
                "Clasificación": "Classification",
                "Regresión": "Regression",
                
                # Cálculo
                "Derivadas": "Derivatives Calculus",
                "Integrales": "Integrals Calculus",
                "Límites": "Limits Calculus",
                "Funciones": "Mathematical Functions"
            }
            
            # Mejorar el prompt
            enhanced_prompt = prompt
            for spanish, english in prompt_translations.items():
                if spanish.lower() in prompt.lower():
                    enhanced_prompt = enhanced_prompt.replace(spanish, english)
            
            # Crear prompt educativo específico
            final_prompt = f"Simple educational illustration about {enhanced_prompt}. Minimalist style, clean white background, professional diagram, no text, suitable for academic presentation"
            
            # Obtener modelo de DALL-E desde env (dall-e-2 o dall-e-3)
            dalle_model = os.getenv("DALLE_MODEL", "dall-e-3")
            
            # Tamaño según modelo
            size = "1024x1024" if dalle_model == "dall-e-3" else "512x512"
            
            print(f"      🎨 Generando: {prompt[:40]}...")
            print(f"      📝 Prompt: {final_prompt[:60]}...")
            print(f"      🤖 Modelo: {dalle_model}")
            
            # Generar imagen
            response = self.openai_client.images.generate(
                model=dalle_model,
                prompt=final_prompt,
                size=size,
                n=1
            )
            
            # Descargar imagen
            image_url = response.data[0].url
            img_data = requests.get(image_url).content
            print(f"      ✅ Imagen generada")
            
            return BytesIO(img_data)
            
        except Exception as e:
            print(f"      ⚠️  Error generando imagen: {str(e)}")
            return None
    
    def _add_title_slide(self, prs: Presentation, title: str, subtitle: str):
        """Agrega slide de título con diseño profesional y moderno"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        shapes = slide.shapes
        
        # Fondo con gradiente azul-morado moderno
        background = shapes.add_shape(
            1,  # Rectangle
            Inches(0), Inches(0),
            Inches(10), Inches(7.5)
        )
        background.fill.solid()
        background.fill.fore_color.rgb = RGBColor(30, 58, 138)  # Azul oscuro
        background.line.fill.background()
        
        # Rectángulo decorativo superior
        accent = shapes.add_shape(
            1,
            Inches(0), Inches(0),
            Inches(10), Inches(1.5)
        )
        accent.fill.solid()
        accent.fill.fore_color.rgb = RGBColor(59, 130, 246)  # Azul brillante
        accent.line.fill.background()
        
        # Título
        left = Inches(0.8)
        top = Inches(2.8)
        width = Inches(8.4)
        height = Inches(2)
        
        title_box = shapes.add_textbox(left, top, width, height)
        title_frame = title_box.text_frame
        title_frame.text = title
        title_frame.word_wrap = True  # Permitir salto de línea
        
        # Formato del título
        title_para = title_frame.paragraphs[0]
        title_para.alignment = PP_ALIGN.CENTER
        title_para.font.size = Pt(40)  # Reducido de 54 a 40
        title_para.font.bold = True
        title_para.line_spacing = 1.2  # Espaciado entre líneas
        title_para.font.color.rgb = RGBColor(255, 255, 255)  # Blanco
        
        # Subtítulo
        subtitle_top = Inches(4.2)
        subtitle_box = shapes.add_textbox(left, subtitle_top, width, Inches(0.8))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = subtitle
        
        subtitle_para = subtitle_frame.paragraphs[0]
        subtitle_para.alignment = PP_ALIGN.CENTER
        subtitle_para.font.size = Pt(28)
        subtitle_para.font.italic = True
        subtitle_para.font.color.rgb = RGBColor(203, 213, 225)  # Gris claro
        
        # Footer
        footer_top = Inches(6.5)
        footer_box = shapes.add_textbox(left, footer_top, width, Inches(0.5))
        footer_frame = footer_box.text_frame
        footer_frame.text = "Generado por Edu-Nexus"
        
        footer_para = footer_frame.paragraphs[0]
        footer_para.alignment = PP_ALIGN.CENTER
        footer_para.font.size = Pt(14)
        footer_para.font.color.rgb = RGBColor(148, 163, 184)
    
    def _add_content_slide(self, prs: Presentation, slide_data: Dict):
        """Agrega slide de contenido con diseño moderno"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
        shapes = slide.shapes
        
        # Fondo blanco limpio
        background = shapes.add_shape(
            1,
            Inches(0), Inches(0),
            Inches(10), Inches(7.5)
        )
        background.fill.solid()
        background.fill.fore_color.rgb = RGBColor(248, 250, 252)  # Gris muy claro
        background.line.fill.background()
        
        # Barra superior con color
        header = shapes.add_shape(
            1,
            Inches(0), Inches(0),
            Inches(10), Inches(1.2)
        )
        header.fill.solid()
        header.fill.fore_color.rgb = RGBColor(59, 130, 246)  # Azul
        header.line.fill.background()
        
        # Título del slide
        title_box = shapes.add_textbox(Inches(0.8), Inches(0.3), Inches(8.4), Inches(0.6))
        title_frame = title_box.text_frame
        title_frame.text = f"📌 {slide_data.get('title', 'Sin título')}"
        
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(36)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(255, 255, 255)  # Blanco sobre azul
        
        # Generar imagen si está habilitado
        image_stream = None
        if self.use_images and self.openai_client:
            slide_title = slide_data.get('title', '')
            image_stream = self._generate_image(slide_title)
        
        # Ajustar layout según si hay imagen o no
        if image_stream:
            # Con imagen: texto a la izquierda, imagen a la derecha
            content_box = shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(4.5), Inches(5))
            
            # Agregar imagen a la derecha
            try:
                shapes.add_picture(
                    image_stream,
                    Inches(5.8), Inches(2.2),
                    width=Inches(3.5)
                )
            except Exception as e:
                print(f"      ⚠️  Error agregando imagen: {str(e)}")
        else:
            # Sin imagen: texto ocupa todo el ancho
            content_box = shapes.add_textbox(Inches(1.2), Inches(1.8), Inches(7.6), Inches(5))
        
        text_frame = content_box.text_frame
        text_frame.word_wrap = True
        
        # Iconos para bullets
        bullet_icons = ["✓", "➤", "●", "◆", "▸"]
        
        bullets = slide_data.get('bullets', [])
        for i, bullet in enumerate(bullets):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            
            # Agregar icono visual
            icon = bullet_icons[i % len(bullet_icons)]
            p.text = f"{icon}  {bullet}"
            p.level = 0
            p.font.size = Pt(20 if image_stream else 22)  # Texto más pequeño si hay imagen
            p.font.bold = False
            p.space_before = Pt(14 if image_stream else 16)
            p.space_after = Pt(6 if image_stream else 8)
            p.font.color.rgb = RGBColor(30, 41, 59)  # Gris oscuro
        
        # Número de slide (footer)
        slide_num = len(prs.slides)
        num_box = shapes.add_textbox(Inches(9), Inches(7), Inches(0.5), Inches(0.3))
        num_frame = num_box.text_frame
        num_frame.text = str(slide_num)
        num_para = num_frame.paragraphs[0]
        num_para.alignment = PP_ALIGN.RIGHT
        num_para.font.size = Pt(12)
        num_para.font.color.rgb = RGBColor(150, 150, 150)
    
    def _add_two_column_slide(self, prs: Presentation, slide_data: Dict):
        """Agrega slide con dos columnas"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        shapes = slide.shapes
        
        # Título
        title_box = shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
        title_frame = title_box.text_frame
        title_frame.text = slide_data.get('title', 'Sin título')
        
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(32)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(31, 78, 120)
        
        # Columna izquierda
        left_box = shapes.add_textbox(Inches(0.5), Inches(2), Inches(4.5), Inches(4.5))
        left_frame = left_box.text_frame
        left_frame.word_wrap = True
        
        left_content = slide_data.get('left_column', [])
        for i, item in enumerate(left_content):
            if i == 0:
                p = left_frame.paragraphs[0]
            else:
                p = left_frame.add_paragraph()
            p.text = item
            p.font.size = Pt(18)
        
        # Columna derecha
        right_box = shapes.add_textbox(Inches(5.5), Inches(2), Inches(4), Inches(4.5))
        right_frame = right_box.text_frame
        right_frame.word_wrap = True
        
        right_content = slide_data.get('right_column', [])
        for i, item in enumerate(right_content):
            if i == 0:
                p = right_frame.paragraphs[0]
            else:
                p = right_frame.add_paragraph()
            p.text = item
            p.font.size = Pt(18)
    
    def _add_summary_slide(self, prs: Presentation, slide_data: Dict):
        """Agrega slide de resumen con diseño impactante"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        shapes = slide.shapes
        
        # Fondo con gradiente
        background = shapes.add_shape(
            1,
            Inches(0), Inches(0),
            Inches(10), Inches(7.5)
        )
        background.fill.solid()
        background.fill.fore_color.rgb = RGBColor(15, 23, 42)  # Azul muy oscuro
        background.line.fill.background()
        
        # Título grande centrado con emoji
        title_box = shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(1.2))
        title_frame = title_box.text_frame
        title_frame.text = f"🎯 {slide_data.get('title', 'Resumen')}"
        
        title_para = title_frame.paragraphs[0]
        title_para.alignment = PP_ALIGN.CENTER
        title_para.font.size = Pt(48)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(255, 255, 255)  # Blanco
        
        # Puntos clave con diseño moderno
        content_box = shapes.add_textbox(Inches(1.5), Inches(3.2), Inches(7), Inches(3.5))
        text_frame = content_box.text_frame
        text_frame.word_wrap = True
        
        bullets = slide_data.get('bullets', [])
        summary_icons = ["⭐", "💡", "🚀", "✨", "🎓"]
        
        for i, bullet in enumerate(bullets):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            
            icon = summary_icons[i % len(summary_icons)]
            p.text = f"{icon}  {bullet}"
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(24)
            p.font.bold = True
            p.space_before = Pt(20)
            p.space_after = Pt(10)
            p.font.color.rgb = RGBColor(255, 255, 255)  # Blanco
    
    def _sanitize_filename(self, filename: str) -> str:
        """Limpia el nombre del archivo para evitar problemas con URLs"""
        import re
        import unicodedata
        from datetime import datetime
        
        # Normalizar unicode y remover acentos
        filename = unicodedata.normalize('NFKD', filename)
        filename = filename.encode('ASCII', 'ignore').decode('ASCII')
        
        # Remover caracteres no válidos
        invalid_chars = '<>:"/\\|?*'
        for char in invalid_chars:
            filename = filename.replace(char, '_')
        
        # Reemplazar espacios con guiones bajos
        filename = filename.replace(' ', '_')
        
        # Remover caracteres especiales restantes
        filename = re.sub(r'[^\w\-_]', '', filename)
        
        # Remover guiones bajos múltiples
        filename = re.sub(r'_+', '_', filename)
        
        # Limitar longitud
        if len(filename) > 40:
            filename = filename[:40]
        
        # Agregar timestamp para unicidad
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        filename = f"{filename}_{timestamp}"
        
        return filename.strip('_')


def parse_llm_response_to_slides(llm_response: str) -> Dict:
    """
    Parsea la respuesta del LLM y la convierte en estructura de slides
    
    Args:
        llm_response: Respuesta del LLM en formato estructurado
        
    Returns:
        Diccionario con title, subtitle y slides
    """
    # Intentar parsear como JSON primero
    try:
        data = json.loads(llm_response)
        return data
    except:
        pass
    
    # Si no es JSON, parsear texto estructurado
    lines = llm_response.strip().split('\n')
    
    title = "Presentación"
    subtitle = ""
    slides = []
    current_slide = None
    
    for line in lines:
        line = line.strip()
        if not line:
            continue
        
        # Detectar título principal
        if line.startswith('# '):
            title = line[2:].strip()
        
        # Detectar subtítulo
        elif line.startswith('## ') and not subtitle:
            subtitle = line[3:].strip()
        
        # Detectar nuevo slide
        elif line.startswith('### '):
            if current_slide:
                slides.append(current_slide)
            current_slide = {
                'type': 'content',
                'title': line[4:].strip(),
                'bullets': []
            }
        
        # Detectar bullets
        elif line.startswith('- ') or line.startswith('* '):
            if current_slide:
                current_slide['bullets'].append(line[2:].strip())
        
        # Detectar números
        elif line[0].isdigit() and '. ' in line:
            if current_slide:
                current_slide['bullets'].append(line.split('. ', 1)[1].strip())
    
    # Agregar último slide
    if current_slide:
        slides.append(current_slide)
    
    return {
        'title': title,
        'subtitle': subtitle,
        'slides': slides
    }


# Función de utilidad para testing
if __name__ == "__main__":
    # Test básico
    generator = PresentationGenerator()
    
    test_data = [
        {
            'type': 'content',
            'title': 'Objetivos de Aprendizaje',
            'bullets': [
                'Comprender el concepto de derivada',
                'Aplicar reglas de derivación',
                'Resolver problemas de optimización',
                'Analizar gráficas de funciones'
            ]
        },
        {
            'type': 'content',
            'title': 'Definición de Derivada',
            'bullets': [
                "f'(x) = lim [f(x+h) - f(x)] / h cuando h→0",
                'Interpretación geométrica: pendiente de la tangente',
                'Interpretación física: velocidad instantánea',
                'Notaciones: f\'(x), dy/dx, Df(x)'
            ]
        },
        {
            'type': 'summary',
            'title': 'Resumen',
            'bullets': [
                'La derivada mide la tasa de cambio',
                'Existen múltiples reglas de derivación',
                'Aplicaciones en optimización y análisis'
            ]
        }
    ]
    
    filepath = generator.create_presentation(
        "Derivadas y Aplicaciones",
        "Unidad 2 - Cálculo Diferencial",
        test_data
    )
    
    print(f"✅ Presentación creada: {filepath}")
